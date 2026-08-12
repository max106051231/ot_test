"""
生成 Semi-Shield ISMS 比賽簡報（含實際 UI 截圖）
輸出：docs/Semi-Shield_ISMS_Competition.pptx
截圖：docs/ppt_assets/*.png
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "Semi-Shield_ISMS_Competition.pptx"
ASSETS = ROOT / "docs" / "ppt_assets"

# 配色
BG = RGBColor(5, 11, 20)
PANEL = RGBColor(10, 21, 36)
ACCENT = RGBColor(78, 200, 240)
ACCENT2 = RGBColor(45, 212, 191)
TEXT = RGBColor(232, 238, 247)
MUTED = RGBColor(138, 160, 184)
WHITE = RGBColor(255, 255, 255)
ORANGE = RGBColor(245, 158, 11)


def _img(name: str) -> Path | None:
    p = ASSETS / name
    return p if p.is_file() else None


def _set_slide_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text, size=18, bold=False,
                 color=TEXT, align=PP_ALIGN.LEFT, font_name="Microsoft JhengHei"):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return box


def _add_title_bar(slide, num: str, title: str, subtitle: str = ""):
    _add_textbox(slide, 0.45, 0.28, 0.9, 0.55, num, size=34, bold=True, color=ACCENT)
    _add_textbox(slide, 1.35, 0.28, 8.2, 0.55, title, size=26, bold=True, color=WHITE)
    if subtitle:
        _add_textbox(slide, 1.35, 0.78, 8.2, 0.35, subtitle, size=13, color=MUTED)
    line = slide.shapes.add_shape(1, Inches(0.45), Inches(1.08), Inches(9.1), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def _bullet_block(slide, items: list[str], left=0.55, top=1.25, width=4.0, size=13):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.8))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT if item else MUTED
        p.font.name = "Microsoft JhengHei"
        p.space_after = Pt(6 if item else 2)
        p.line_spacing = 1.15


def _add_image_card(slide, image_path: Path, left, top, width, height, caption: str = ""):
    """圓角感邊框 + 截圖 + 說明文字"""
    frame = slide.shapes.add_shape(1, Inches(left - 0.04), Inches(top - 0.04),
                                   Inches(width + 0.08), Inches(height + 0.08))
    frame.fill.solid()
    frame.fill.fore_color.rgb = PANEL
    frame.line.color.rgb = ACCENT
    frame.line.width = Pt(1.2)
    slide.shapes.add_picture(
        str(image_path),
        Inches(left), Inches(top),
        width=Inches(width), height=Inches(height),
    )
    if caption:
        cap = slide.shapes.add_shape(1, Inches(left), Inches(top + height - 0.02),
                                     Inches(width), Inches(0.42))
        cap.fill.solid()
        cap.fill.fore_color.rgb = RGBColor(8, 16, 28)
        cap.line.fill.background()
        _add_textbox(slide, left + 0.08, top + height + 0.02, width - 0.16, 0.38,
                     caption, size=10, color=MUTED)


def _slide_text_image(slide, num, title, subtitle, bullets, img_name, caption="",
                      text_w=3.85, img_left=4.55, img_w=5.0, img_h=5.35):
    _add_title_bar(slide, num, title, subtitle)
    _bullet_block(slide, bullets, width=text_w)
    p = _img(img_name)
    if p:
        _add_image_card(slide, p, img_left, 1.22, img_w, img_h, caption)


def _three_col_slide(slide, cols: list[tuple[str, str, str]], top=1.35):
    w = 2.95
    for i, (head, sub, body) in enumerate(cols):
        x = 0.45 + i * 3.12
        rect = slide.shapes.add_shape(1, Inches(x), Inches(top), Inches(w), Inches(4.85))
        rect.fill.solid()
        rect.fill.fore_color.rgb = PANEL
        rect.line.color.rgb = ACCENT
        _add_textbox(slide, x + 0.12, top + 0.12, w - 0.24, 0.45, head, size=18, bold=True, color=ACCENT2)
        _add_textbox(slide, x + 0.12, top + 0.58, w - 0.24, 0.35, sub, size=12, bold=True, color=ORANGE)
        _add_textbox(slide, x + 0.12, top + 0.95, w - 0.24, 3.6, body, size=11, color=TEXT)


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ===== 1 封面 =====
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s)
    # 左側色條
    bar = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.12), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT2
    bar.line.fill.background()
    _add_textbox(s, 0.55, 0.9, 5.2, 1.4,
                 "Semi-Shield ISMS\n地端 OT × ISO 27001\n合規 Agent 平台",
                 size=30, bold=True, color=WHITE)
    _add_textbox(s, 0.55, 2.45, 5.0, 0.7,
                 "Wazuh + OpenCTI + Local LLM\n+ Agentic RAG + Evidence Traceability",
                 size=14, color=ACCENT)
    _add_textbox(s, 0.55, 3.35, 5.0, 1.8,
                 "指導老師：＿＿＿＿＿＿ 老師\n"
                 "企業業師：毛勁豪 老師（SenseL 平台）\n"
                 "組員：王昱勻、陳昭宇、鍾世荃\n"
                 "隊名：Semi-Shield",
                 size=13, color=MUTED)
    _add_textbox(s, 0.55, 6.35, 5.5, 0.35,
                 "No Evidence, No Compliance Claim",
                 size=12, bold=True, color=ACCENT2)
    p = _img("03_split_monitor_chat.png")
    if p:
        _add_image_card(s, p, 5.35, 0.75, 4.35, 5.85,
                        "實際平台：監控戰情 × AI Agent 並排檢視")

    # ===== 2 大綱 =====
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s)
    _add_textbox(s, 0.55, 0.45, 4, 0.55, "OUTLINE", size=30, bold=True, color=ACCENT)
    cols_l = ["01 Problem", "02 Customer", "03 UVP", "04 Architecture"]
    cols_r = ["05 Technical", "06 Demo", "07 Validation", "08 Roadmap"]
    for i, (a, b) in enumerate(zip(cols_l, cols_r)):
        y = 1.35 + i * 0.72
        _add_textbox(s, 0.7, y, 4.2, 0.5, a, size=17, color=TEXT)
        _add_textbox(s, 5.2, y, 4.2, 0.5, b, size=17, color=TEXT)
    p = _img("01_monitor_dashboard.png")
    if p:
        _add_image_card(s, p, 5.0, 4.35, 4.5, 2.55, "ISO 27001 監控戰情室（實機截圖）")

    # ===== 3 Problem =====
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s)
    _add_title_bar(s, "01", "關鍵問題：SOC 與合規的斷層", "Problem")
    _three_col_slide(s, [
        ("告警疲勞", "SOC 團隊的惡夢",
         "Wazuh / CTI / IR 每日大量事件\n80% 時間 triage 誤報\n真實威脅淹沒雜訊\n→ 事件 ≠ 稽核證據"),
        ("合規斷層", "No Evidence, No Claim",
         "ISO 27001 需控制項對應\n+ 可追溯 evidence\n人工整理報告 4 週+\n→ 顧問成本高"),
        ("AI 報告幻覺", "傳統 LLM 局限",
         "單次生成、無 evidence_id\n易 hallucination\n無 Human Review\n→ 不能當 ISMS 證據"),
    ])

    # ===== 4 Customer =====
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s)
    _add_title_bar(s, "02", "目標客群", "Customer Segments")
    _bullet_block(s, [
        "【主要產業】",
        "• 半導體／電子製造（Cisco Catalyst 管理面、無塵室 OT）",
        "• 智慧製造（ISO 27001 + OT 日誌合規證據）",
        "• 關鍵基礎設施（電力、水務、石化）",
        "",
        "【客戶特徵】",
        "• 已有 SOC，Annex A 證據整理靠人工",
        "• 敏感資料不能上雲 → Local LLM",
        "• 合規壓力：27001、IEC 62443、NIS2",
    ], width=8.5, size=15)

    # ===== 5 UVP =====
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s)
    _add_title_bar(s, "03", "研究目的與創新性", "UVP")
    _bullet_block(s, [
        "【目的】SOC 事件 → ISO 27001 可審查證據（4 週 → ≤2 天）",
        "",
        "1. evidence_id + hash chain（No Evidence, No Claim）",
        "2. Agent：Collector → Auditor → Reviewer → Reporter",
        "3. Local LLM 地端（Qwen3-4B / Gemma LoRA）",
        "4. Agentic RAG（681 筆 OT/ISO 知識庫）",
        "5. Human Review Guardrails（RoBERTa + 覆核佇列）",
        "6. syslog → Annex A 六控制項自動對映",
    ], width=8.8, size=14)

    # ===== 6 Architecture + 截圖 =====
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s)
    _slide_text_image(
        s, "04", "系統架構", "Five-Layer Architecture",
        [
            "⑤ Web Console",
            "  platform / monitor / chat / compliance",
            "④ Guardrails + Human Review",
            "③ ISO 27001 Mapping Matrix",
            "② RAG & Local LLM Reasoning",
            "① Agent Workflow + Data",
            "  Wazuh · OpenCTI · evidence_service",
        ],
        "03_split_monitor_chat.png",
        "整合平台：監控 + AI 對話並排（RAG ON · 離線就緒）",
        text_w=3.5, img_left=4.15, img_w=5.45, img_h=5.5,
    )

    # ===== 7 Tech Agent =====
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s)
    _add_title_bar(s, "05", "技術說明 (1/4)", "Agent Workflow")
    _bullet_block(s, [
        "Collector → scan OT / 註冊 evidence_id",
        "Auditor → RAG + LLM 控制項診斷",
        "Reviewer → Guardrails + review_queue",
        "Reporter → PDF/TXT 合規報告",
        "",
        "API：POST /api/compliance/pipeline",
        "Phase 2：LangGraph / Wazuh / OpenCTI",
    ], width=8.5, size=15)

    # ===== 8 ISO Mapping + 監控截圖 =====
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s)
    _slide_text_image(
        s, "05", "技術說明 (2/4)", "ISO 27001 OT 控制項對映",
        [
            "Cisco syslog → 6 控制項：",
            "• A.5.15 存取控制",
            "• A.5.19 供應鏈",
            "• A.8.7 惡意程式",
            "• A.8.8 弱點修補",
            "• A.8.19 組態變更",
            "• A.8.24 傳輸加密",
            "",
            "Phase 1：71 項矩陣 / 目標 ≥80%",
        ],
        "01_monitor_dashboard.png",
        "監控戰情：六控制項 KPI + 不符合統計 + 日誌串流",
        text_w=3.15, img_left=3.85, img_w=5.75, img_h=5.45,
    )

    # ===== 9 Evidence + 診斷截圖 =====
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s)
    _slide_text_image(
        s, "05", "技術說明 (3/4)", "Guardrails + Evidence 追溯",
        [
            "Guardrails 三層：",
            "  規則 → allowlist → RoBERTa",
            "  borderline → Human Review",
            "",
            "evidence_id：",
            "  EV-{site}-{ctrl}-{UTC}-{hash}",
            "  SHA-256 hash chain",
            "",
            "LLM 診斷 → 三段落合規報告",
            "（摘要 / 風險 / 修補）",
        ],
        "04_ai_diagnosis_modal.png",
        "A.8.19 地端 LLM 診斷報告（實際執行畫面）",
        text_w=3.2, img_left=3.95, img_w=5.65, img_h=5.45,
    )

    # ===== 10 Local LLM + 模型選單 =====
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s)
    _slide_text_image(
        s, "05", "技術說明 (4/4)", "Local LLM 微調與熱切換",
        [
            "Production：",
            "  Qwen3-4B LoRA (qwen3_4b_ot)",
            "Alternatives：",
            "  Gemma-2-2b-it LoRA",
            "  Phi-4-mini / Qwen2.5-3B",
            "",
            "349 筆 OT/ISO 訓練資料",
            "CUDA 4-bit QLoRA · CPU failover",
            "Edge：run_edge.bat",
        ],
        "02_model_selector.png",
        "微調前 / 微調後模型熱切換（gemma_2b_ot 等）",
        text_w=3.2, img_left=3.95, img_w=5.65, img_h=5.45,
    )

    # ===== 11 Demo 圖表 =====
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s)
    _add_title_bar(s, "06", "成果展示 (1/2)", "AI Agent 圖表與合規視覺化")
    p1 = _img("05_chat_charts.png")
    p2 = _img("06_platform_charts.png")
    if p1:
        _add_image_card(s, p1, 0.45, 1.2, 4.55, 5.55,
                        "聊天 Agent：OT 控制項事件量長條圖 + 圓餅圖")
    if p2:
        _add_image_card(s, p2, 5.15, 1.2, 4.45, 5.55,
                        "gemma_2b_ot：合規控制項圖表生成")
    _add_textbox(s, 0.45, 6.85, 9.1, 0.35,
                 "Demo：http://localhost:2000/platform ｜ RAG ON(681) ｜ 模型可切換",
                 size=11, color=MUTED, align=PP_ALIGN.CENTER)

    # ===== 12 Demo 全景 =====
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s)
    _add_title_bar(s, "06", "成果展示 (2/2)", "監控 × 對話 一站式合規作業")
    p = _img("03_split_monitor_chat.png")
    if p:
        _add_image_card(s, p, 0.55, 1.18, 8.9, 5.65,
                        "左：ISO 27001 監控戰情 ｜ 右：Semi-Shield Cyber Agent + 對話紀錄")
    _bullet_block(s, [
        "1. syslog 掃描 → metrics",
        "2. 單項 AI 診斷 → evidence",
        "3. 聊天 + RAG → 圖表",
        "4. PDF/TXT 匯出",
    ], left=0.55, top=6.72, width=9.0, size=11)

    # ===== 13 KPI =====
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s)
    _add_title_bar(s, "07", "成果驗證 (1/2)", "量化 KPI 目標")
    rows = [
        ("Annex A 覆蓋率", "≥ 80%"),
        ("控制項對映準確率", "≥ 85%"),
        ("Human review 採納率", "≥ 75%"),
        ("Evidence 可追溯率", "100%"),
        ("報告生成時間", "≤ 2 天（基線 4 週）"),
        ("合規成本降低", "≥ 70%"),
        ("RAG 知識庫", "681 筆"),
        ("OT 自動監控", "6 項 → Phase 2: 15+"),
    ]
    y = 1.42
    for k, v in rows:
        bg = s.shapes.add_shape(1, Inches(0.55), Inches(y - 0.04), Inches(8.9), Inches(0.48))
        bg.fill.solid()
        bg.fill.fore_color.rgb = PANEL if (int(y * 10) % 2 == 0) else BG
        bg.line.fill.background()
        _add_textbox(s, 0.7, y, 3.5, 0.35, k, size=13, bold=True, color=ACCENT2)
        _add_textbox(s, 4.3, y, 4.8, 0.35, v, size=13, color=TEXT)
        y += 0.56

    # ===== 14 貢獻 =====
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s)
    _add_title_bar(s, "07", "成果驗證 (2/2)", "產業貢獻")
    _three_col_slide(s, [
        ("合規自動化", "SOC → ISMS",
         "事件轉 evidence\n報告 4週→2天\n降低顧問依賴"),
        ("地端 AI", "Local LLM",
         "資料不離廠\nHuman Review\n優於 SaaS GRC"),
        ("半導體 OT", "SenseL anchoring",
         "Cisco syslog\nIEC 62443 對照\nEdge 可部署"),
    ], top=1.45)

    # ===== 15 Roadmap =====
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s)
    _add_title_bar(s, "08", "未來擴充性", "Roadmap")
    _bullet_block(s, [
        "Wazuh + OpenCTI 原生整合 · CTI feedback loop",
        "Annex A 6 → 15+ 自動監控",
        "LangGraph / CrewAI Agent 编排",
        "NeMo Guardrails + LlamaGuard",
        "Edge：樹莓 Pi · MQTT / OPC UA",
    ], width=8.5, size=15)

    # ===== 16 Thank you =====
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s)
    _add_textbox(s, 0.8, 2.0, 8.4, 1.6, "THANK YOU\nFOR YOUR\nATTENTION",
                 size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(s, 0.8, 4.85, 8.4, 0.45,
                 "Semi-Shield ISMS ｜ No Evidence, No Compliance Claim",
                 size=15, color=ACCENT, align=PP_ALIGN.CENTER)
    p = _img("01_monitor_dashboard.png")
    if p:
        _add_image_card(s, p, 2.5, 5.45, 5.0, 1.65, "semi-shield.local:2000/platform")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print("Saved:", str(OUT))
    print("Assets:", ASSETS)


if __name__ == "__main__":
    build()
